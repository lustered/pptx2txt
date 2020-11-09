#!/usr/bin/env python3
from pptx import Presentation
import glob
from tika import parser


def grab_text():
    for ppt in glob.glob("*.pptx"):
        prs = Presentation(ppt)
        f = open(''.join([ppt.strip('.pptx'), '.txt']),
                 'w', encoding='utf-8', errors='ignore')
        print("Processed file: ", ppt)
        print("----------------------")
        for index, slide in enumerate(prs.slides):
            # f.write("".join(['\n\n', "Slide #: ", str(index), '\n']))
            f.write('\n\n')
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = bytes(shape.text, 'UTF-8').decode(
                        'UTF-8', 'ignore')
                    f.write(str(text))


def using_tika(filename):
    parsed = parser.from_file(filename)
    # print(parsed["metadata"])
    print(parsed["content"].strip(' '))
    f = open(''.join([filename.strip('.pptx'), '.txt']),
             'w', encoding='utf-8', errors='ignore')
    f.write(parsed["content"].strip(' '))
    f.close()


if __name__ == "__main__":
    # FILENAME = 'Diabetes PPT-1.pptx'
    # using_tika(FILENAME)
    grab_text()
